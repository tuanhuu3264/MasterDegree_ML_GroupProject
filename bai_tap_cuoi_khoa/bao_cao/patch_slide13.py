# -*- coding: utf-8 -*-
"""Slide 13: thay textbox importance = BIEU DO COT chung minh 2 bien nhiet do la thu pham."""
import sys; sys.stdout.reconfigure(encoding='utf-8')
import numpy as np, pandas as pd, warnings; warnings.filterwarnings('ignore')
import matplotlib; matplotlib.use('Agg'); import matplotlib.pyplot as plt
from matplotlib import font_manager
from sklearn.ensemble import RandomForestClassifier
from pptx import Presentation
from pptx.util import Inches, Emu

for fn in ['Segoe UI','Arial','DejaVu Sans']:
    try: font_manager.findfont(fn,fallback_to_default=False); plt.rcParams['font.family']=fn; break
    except: pass
plt.rcParams['figure.dpi']=150

# --- 1) tinh lai importance cua drift classifier (A=0 / B=1) ---
train=pd.read_csv('Data_Final/train.csv'); test=pd.read_csv('Data_Final/test.csv')
NUM=['nhiet_do_moi_truong','nhiet_do_quy_trinh','toc_do_quay','momen_xoan','do_mon_dao']
X=pd.concat([train[NUM],test[NUM]],ignore_index=True)
y=np.r_[np.zeros(len(train)),np.ones(len(test))]
rf=RandomForestClassifier(n_estimators=300,random_state=42,n_jobs=-1,class_weight='balanced')
rf.fit(X,y)
imp=pd.Series(rf.feature_importances_,index=NUM).sort_values()
name={'nhiet_do_moi_truong':'nhiệt độ môi trường','nhiet_do_quy_trinh':'nhiệt độ máy',
      'toc_do_quay':'tốc độ quay','momen_xoan':'lực xoắn','do_mon_dao':'độ mòn dao'}
is_temp={'nhiet_do_moi_truong':True,'nhiet_do_quy_trinh':True}
temp_total=sum(imp[k] for k in imp.index if is_temp.get(k,False))*100

# --- 2) ve bar ngang, 2 bien nhiet do mau do ---
fig,ax=plt.subplots(figsize=(8.7,2.0))
colors=['#dc2626' if is_temp.get(k,False) else '#94a3b8' for k in imp.index]
bars=ax.barh(range(len(imp)),imp.values*100,color=colors,zorder=3,height=0.66)
ax.set_yticks(range(len(imp))); ax.set_yticklabels([name[k] for k in imp.index],fontsize=10.5)
for i,v in enumerate(imp.values*100):
    ax.text(v+0.6,i,f'{v:.0f}%',va='center',fontsize=10,fontweight='bold')
ax.set_xlim(0,40); ax.set_xlabel('Độ quan trọng khi đoán A/B (%)',fontsize=9.5)
for s in ['top','right']: ax.spines[s].set_visible(False)
ax.grid(axis='x',alpha=0.25,zorder=0)
ax.set_title(f'Feature importance của Drift Classifier — 2 biến nhiệt độ (đỏ) chiếm ~{temp_total:.0f}% → thủ phạm chính',
             fontsize=10.5,fontweight='bold',pad=6)
plt.tight_layout()
IMG='bao_cao/assets/10_drift_importance.png'
plt.savefig(IMG,bbox_inches='tight'); plt.close()
print(f'saved {IMG}  (2 nhiet do = {temp_total:.1f}%)')

# --- 3) va slide 13: xoa textbox importance cu, chen bieu do ---
F='bao_cao/SLIDE_THUYET_TRINH_v6.pptx'
prs=Presentation(F); slide=prs.slides[12]
for sh in list(slide.shapes):
    if sh.has_text_frame and sh.text_frame.text.startswith('Feature importance'):
        sh._element.getparent().remove(sh._element); print('removed old importance textbox')
slide.shapes.add_picture(IMG, Inches(0.62), Inches(5.62), width=Inches(7.30))
prs.save(F)
print('Slide 13 patched with bar chart.')

# verify
p2=Presentation(F); s=p2.slides[12]
for sh in s.shapes:
    if sh.left is None: continue
    L,T,W,H=Emu(sh.left).inches,Emu(sh.top).inches,Emu(sh.width).inches,Emu(sh.height).inches
    tag='PIC' if sh.shape_type==13 else (repr(sh.text_frame.text[:28]) if sh.has_text_frame and sh.text_frame.text.strip() else sh.name)
    print(f'  T={T:.2f} B={T+H:.2f} R={L+W:.2f}  {tag}')
