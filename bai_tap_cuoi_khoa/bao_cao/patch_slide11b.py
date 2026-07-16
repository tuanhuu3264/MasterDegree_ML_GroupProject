# -*- coding: utf-8 -*-
"""Slide 11 tinh chinh: thu nho font cot Cong thuc/Y nghia, day 2 textbox duoi xuong."""
import sys; sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt

F = 'bao_cao/SLIDE_THUYET_TRINH_v6.pptx'
prs = Presentation(F)
slide = prs.slides[10]

# 1) thu nho font cot 1 (Cong thuc) va 2 (Y nghia), hang du lieu 1..3
tbl = next(sh for sh in slide.shapes if sh.has_table).table
for r in range(1, len(tbl.rows)):
    for c in (1, 2):
        cell = tbl.cell(r, c)
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(11)

# 2) day textbox 'Vi sao hay' va 'Ghi chu' xuong de tranh cham day bang
for sh in slide.shapes:
    if not (sh.has_text_frame and sh.left is not None): continue
    t = sh.text_frame.text
    if t.startswith('Vì sao hay'):
        sh.top = Inches(4.45)
    elif t.startswith('Ghi chú'):
        sh.top = Inches(5.70)

prs.save(F)
print('Slide 11: da thu nho font + day textbox xuong.')

# verify geometry
from pptx.util import Emu
p2 = Presentation(F); s = p2.slides[10]
for sh in s.shapes:
    if sh.left is None: continue
    L,T,W,H = Emu(sh.left).inches,Emu(sh.top).inches,Emu(sh.width).inches,Emu(sh.height).inches
    tag='TABLE' if sh.has_table else (sh.text_frame.text[:28] if sh.has_text_frame and sh.text_frame.text.strip() else sh.name)
    print(f'  T={T:.2f} B={T+H:.2f}  {tag}')
