# -*- coding: utf-8 -*-
"""Va slide 12 (them cot KS-D) va slide 13 (them feature importance drift)."""
import sys, copy; sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
SRC = 'bao_cao/SLIDE_THUYET_TRINH_v5.pptx'
DST = 'bao_cao/SLIDE_THUYET_TRINH_v6.pptx'

prs = Presentation(SRC)

# ============ SLIDE 12 (index 11): them cot KS-D giua PSI va Phan loai ============
slide12 = prs.slides[11]
tbl_shape = next(sh for sh in slide12.shapes if sh.has_table)
tbl = tbl_shape.table
tblEl = tbl._tbl

# 1) them gridCol tai vi tri 2 (sau PSI)
tblGrid = tblEl.find(A+'tblGrid')
gridCols = tblGrid.findall(A+'gridCol')
newcol = copy.deepcopy(gridCols[1])          # copy cot PSI de thua width/format
gridCols[1].addnext(newcol)

# 2) moi hang: chen tc moi (copy tc PSI) tai vi tri 2, set text KS
ks_vals = ['KS-D', '0,43', '0,31', '0,17 / 0,14', '≈ 0,01', '0,009']
rows = tblEl.findall(A+'tr')
for ri, tr in enumerate(rows):
    tcs = tr.findall(A+'tc')
    newtc = copy.deepcopy(tcs[1])            # copy o PSI cung hang
    ts = newtc.findall('.//'+A+'t')
    if ts:
        ts[0].text = ks_vals[ri]
        for extra in ts[1:]:
            extra.text = ''
    tcs[1].addnext(newtc)

# 3) can lai do rong 4 cot cho vua (tong ~6.4in, truoc anh o L=7.10)
widths = [Inches(2.30), Inches(0.95), Inches(1.15), Inches(1.55)]
for c, w in zip(tbl.columns, widths):
    c.width = w

print('Slide 12: da them cot KS-D ->', len(tbl.columns), 'cot')

# ============ SLIDE 13 (index 12): them feature importance drift ============
slide13 = prs.slides[12]
tb = slide13.shapes.add_textbox(Inches(0.62), Inches(5.75), Inches(7.30), Inches(1.35))
tf = tb.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
r1 = p1.add_run(); r1.text = 'Feature importance (mô hình đoán A/B):'
r1.font.bold = True; r1.font.size = Pt(13); r1.font.color.rgb = RGBColor(0x14,0x24,0x3d)

p2 = tf.add_paragraph()
r2 = p2.add_run()
r2.text = ('nhiệt độ MT 33%  ·  nhiệt độ máy 21%  ·  tốc độ 16%  ·  '
           'lực xoắn 16%  ·  độ mòn dao 14%')
r2.font.size = Pt(12.5); r2.font.color.rgb = RGBColor(0x22,0x22,0x22)

p3 = tf.add_paragraph()
r3 = p3.add_run()
r3.text = '⇒ 2 biến nhiệt độ chiếm ~54% → thủ phạm chính gây shift.'
r3.font.size = Pt(12.5); r3.font.italic = True; r3.font.color.rgb = RGBColor(0xb0,0x1a,0x1a)

print('Slide 13: da them textbox feature importance')

prs.save(DST)
print('SAVED ->', DST)

# ---- verify ----
prs2 = Presentation(DST)
t2 = next(sh for sh in prs2.slides[11].shapes if sh.has_table).table
print('\nVERIFY slide 12 table:')
for r in range(len(t2.rows)):
    print('  | ' + ' | '.join(t2.cell(r,c).text for c in range(len(t2.columns))))
