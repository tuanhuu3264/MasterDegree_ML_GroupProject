# -*- coding: utf-8 -*-
"""Sua claim sai 'toc_do_quay hai dau' -> toc do dau THAP; momen moi la hai dau.
Sua slide 8 + notebook cell 28,29."""
import sys, json; sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

# ========== 1) SLIDE 8 ==========
F='bao_cao/SLIDE_THUYET_TRINH_v6.pptx'
prs=Presentation(F)
sh=[s for s in prs.slides[7].shapes if s.has_text_frame and s.text_frame.text.startswith('1.3')][0]
paras=sh.text_frame.paragraphs
# giu 5 para; sua noi dung para 1,2,3
new_lines={
    1:'• Độ mòn dao rõ nhất: máy hỏng ~177 vs lành ~123 → dồn mức mòn CAO.',
    2:'• Biến khác tương quan tuyến tính ≈ 0 nhưng KHÔNG vô dụng:',
    3:'• Tốc độ quay dồn đuôi THẤP; mômen xoắn dồn CẢ HAI ĐẦU → quá/thiếu tải.',
}
for idx,txt in new_lines.items():
    p=paras[idx]
    if p.runs:
        p.runs[0].text=txt
        for extra in p.runs[1:]:
            extra.text=''
prs.save(F)
print('Slide 8: da sua bullet 1.3')
for p in Presentation(F).slides[7].shapes:
    if p.has_text_frame and p.text_frame.text.startswith('1.3'):
        for ln in p.text_frame.text.split(chr(10)): print('   ', ln)

# ========== 2) NOTEBOOK cell 28, 29 ==========
NB='bai_tap_cuoi_khoa.ipynb'
nb=json.load(open(NB,encoding='utf-8'))

# cell 28 (markdown 💡)
c28=('> 💡 **Hiểu đơn giản.** Xem máy hỏng thường "nằm ở vùng nào" của mỗi biến.\n'
     '>\n'
     '> 🔑 **Rút ra:** Máy hỏng có **mômen xoắn** dồn về **HAI đầu** (rất cao HOẶC rất thấp), '
     'còn **tốc độ quay** dồn về **đuôi THẤP** (quay chậm) → hỏng do quá tải/thiếu tải, '
     'không phải một ngưỡng đơn giản.')
nb['cells'][28]['source']=c28.splitlines(keepends=True)

# cell 29 (markdown Nhan xet)
c29=('**Nhận xét (a).** `do_mon_dao` của máy hỏng lệch hẳn sang **vùng mòn cao** (đuôi phải).\n'
     'Đáng chú ý: `momen_xoan` của máy hỏng có **hai cụm ở hai đuôi** (rất thấp *hoặc* rất cao),\n'
     'còn `toc_do_quay` dồn về **đuôi thấp** (quay chậm) — dấu hiệu **hỏng do quá tải/thiếu tải**\n'
     'phụ thuộc *tổ hợp* biến, không phải một ngưỡng đơn. Chính kiểu tín hiệu **phi tuyến** này\n'
     'khiến LogReg tuyến tính bó tay và mô hình cây sẽ thắng.')
nb['cells'][29]['source']=c29.splitlines(keepends=True)

json.dump(nb, open(NB,'w',encoding='utf-8'), ensure_ascii=False, indent=1)
print('\nNotebook: da sua cell 28, 29')
print('  cell28:', ''.join(nb['cells'][28]['source']).split('Rút ra:')[1][:120].strip())
print('  cell29:', ''.join(nb['cells'][29]['source'])[:120].strip())
