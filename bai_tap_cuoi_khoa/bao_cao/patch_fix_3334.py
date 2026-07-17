# -*- coding: utf-8 -*-
"""Sua: (1) tieu de slide 15/16 de phan biet; (2) so reweighting sai o slide 15 & 18."""
import sys; sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

F='bao_cao/SLIDE_THUYET_TRINH_v6.pptx'
prs=Presentation(F)

def set_run(shape, old, new):
    """Thay text 1 run khop 'old' (chinh xac) -> 'new'. Tra ve True neu thay."""
    if not shape.has_text_frame: return False
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.text == old:
                r.text = new; return True
    return False

log=[]
# ---- SLIDE 15 (index 14) ----
s15=prs.slides[14]
for sh in s15.shapes:
    if set_run(sh,'Phần 3.3 — Xử lý lệch #1: chú ý các máy "giống B"',
                  'Phần 3.3 — Xử lý lệch bằng TRỌNG SỐ (Reweighting)'):
        log.append('S15 title')
    if set_run(sh,'Kết quả trước → sau: AUC-PR chỉ nhích nhẹ 0,669 → 0,672.',
                  'Kết quả trước → sau: AUC-PR chỉ nhích nhẹ 0,672 → 0,674.'):
        log.append('S15 bullet AUC-PR')
    if set_run(sh,'AUC-PR:  0,669  →  0,672','AUC-PR:  0,672  →  0,674'):
        log.append('S15 box AUC-PR')
    if set_run(sh,'F1:          0,770  →  0,764','F1:          0,769  →  0,772'):
        log.append('S15 box F1')

# ---- SLIDE 16 (index 15) ----
s16=prs.slides[15]
for sh in s16.shapes:
    if set_run(sh,'Phần 3.4 — Xử lý lệch #2: chọn ngưỡng & chấm thử trung thực',
                  'Phần 3.4 — Chọn NGƯỠNG trung thực (IWV)'):
        log.append('S16 title')

# ---- SLIDE 18 (index 17) table: sua AUC-ROC & AUC-PR 2 dong v3/v4 ----
tbl=next(sh for sh in prs.slides[17].shapes if sh.has_table).table
def fix_cell(r,c,old,new):
    cell=tbl.cell(r,c)
    for p in cell.text_frame.paragraphs:
        for run in p.runs:
            if run.text==old: run.text=new; return True
    return False
# row6 = + reweighting ; row7 = + chỉnh ngưỡng
for row in (6,7):
    if fix_cell(row,1,'0,868','0,866'): log.append(f'S18 r{row} AUC-ROC')
    if fix_cell(row,2,'0,672','0,674'): log.append(f'S18 r{row} AUC-PR')
fix_cell(7,3,'0,774','0,773') and log.append('S18 r7 F1')

prs.save(F)
print('DA SUA:', log)

# verify
p2=Presentation(F)
print('\n-- S15 title:', p2.slides[14].shapes[0].text_frame.text)
print('-- S16 title:', p2.slides[15].shapes[0].text_frame.text)
for sh in p2.slides[14].shapes:
    if sh.has_text_frame and 'AUC-PR:' in sh.text_frame.text:
        print('-- S15 box:', sh.text_frame.text.replace(chr(10),' | '))
t=next(sh for sh in p2.slides[17].shapes if sh.has_table).table
for r in (6,7):
    print(f'-- S18 row{r}:', [t.cell(r,c).text for c in range(len(t.columns))])
