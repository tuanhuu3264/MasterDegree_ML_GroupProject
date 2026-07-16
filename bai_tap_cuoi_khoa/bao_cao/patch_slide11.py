# -*- coding: utf-8 -*-
"""Slide 11: chen cot 'Cong thuc' vao bang 3 dac trung + them ghi chu P/DeltaT."""
import sys, copy; sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
F = 'bao_cao/SLIDE_THUYET_TRINH_v6.pptx'
prs = Presentation(F)
slide = prs.slides[10]  # slide 11

tbl_shape = next(sh for sh in slide.shapes if sh.has_table)
tbl = tbl_shape.table
tblEl = tbl._tbl

# 1) gridCol moi tai vi tri 1 (sau ten dac trung)
tblGrid = tblEl.find(A+'tblGrid')
gridCols = tblGrid.findall(A+'gridCol')
newcol = copy.deepcopy(gridCols[1])       # copy cot 'Y nghia' de thua format text thuong
gridCols[0].addnext(newcol)

# 2) moi hang: chen tc moi (copy o 'Y nghia') tai vi tri 1, set cong thuc
formulas = [
    'Công thức',
    'max(8,6 − ΔT, 0) × max(1380 − tốc, 0)',
    'max(3500 − P, 0) + max(P − 9000, 0)',
    'mòn × mômen − ngưỡng(loại SP)',
]
rows = tblEl.findall(A+'tr')
for ri, tr in enumerate(rows):
    tcs = tr.findall(A+'tc')
    newtc = copy.deepcopy(tcs[1])
    ts = newtc.findall('.//'+A+'t')
    if ts:
        ts[0].text = formulas[ri]
        for extra in ts[1:]:
            extra.text = ''
    tcs[0].addnext(newtc)

# 3) can lai do rong 4 cot (tong = 12,1in)
for c, w in zip(tbl.columns, [Inches(2.35), Inches(4.05), Inches(3.30), Inches(2.40)]):
    c.width = w

# 4) ghi chu dinh nghia P va DeltaT (dat ngay duoi text 'Vi sao hay')
note = slide.shapes.add_textbox(Inches(0.63), Inches(5.35), Inches(12.10), Inches(0.55))
tf = note.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; r = p.add_run()
r.text = ('Ghi chú: ΔT = nhiệt độ máy − nhiệt độ môi trường (K);  '
          'P (công suất cơ) = mômen × tốc độ × 2π/60 (W).  '
          'Hàm max(·,0) = "bản lề": chỉ bật khi máy đi vào vùng nguy.')
r.font.size = Pt(11.5); r.font.italic = True; r.font.color.rgb = RGBColor(0x33,0x33,0x33)

prs.save(F)
print('Slide 11 -> them cot Cong thuc + ghi chu. Cols =', len(tbl.columns))

# verify
p2 = Presentation(F)
t2 = next(sh for sh in p2.slides[10].shapes if sh.has_table).table
for rr in range(len(t2.rows)):
    print('  | ' + ' | '.join(t2.cell(rr,c).text for c in range(len(t2.columns))))
