# -*- coding: utf-8 -*-
"""Slide 8: them bieu do PHAN BO THEO NHAN (ho tro nhan xet 1.3) vao vung trong duoi-trai.
Giu heatmap tuong quan (rubric yeu cau) + thu font text 1.3 cho gon."""
import sys; sys.stdout.reconfigure(encoding='utf-8')
import numpy as np, pandas as pd, warnings; warnings.filterwarnings('ignore')
import matplotlib; matplotlib.use('Agg'); import matplotlib.pyplot as plt
from matplotlib import font_manager
from pptx import Presentation
from pptx.util import Inches, Pt

for fn in ['Segoe UI','Arial','DejaVu Sans']:
    try: font_manager.findfont(fn,fallback_to_default=False); plt.rcParams['font.family']=fn; break
    except: pass
plt.rcParams['figure.dpi']=150

# --- 1) sinh hinh phan bo theo nhan ---
train=pd.read_csv('Data_Final/train.csv'); y=train['hong_hoc']
fail=train[y==1]; ok=train[y==0]
cols=[('do_mon_dao','Độ mòn dao'),('toc_do_quay','Tốc độ quay'),('momen_xoan','Mômen xoắn')]
fig,axes=plt.subplots(1,3,figsize=(9.4,2.75))
for ax,(c,t) in zip(axes,cols):
    bins=np.linspace(train[c].min(),train[c].max(),34)
    ax.hist(ok[c],bins=bins,density=True,alpha=0.5,color='#2563eb',label='không hỏng')
    ax.hist(fail[c],bins=bins,density=True,alpha=0.55,color='#dc2626',label='hỏng')
    ax.set_title(t,fontsize=11); ax.set_yticks([]); ax.grid(alpha=0.2)
    for s in ['top','right']: ax.spines[s].set_visible(False)
axes[0].legend(fontsize=8.5,loc='upper left',framealpha=0.9)
fig.suptitle('Phân bố số đo theo nhãn (Train) — đỏ = máy hỏng',fontsize=11,fontweight='bold',y=1.03)
plt.tight_layout()
IMG='bao_cao/assets/09_density_by_label.png'
plt.savefig(IMG,bbox_inches='tight'); plt.close()
print('saved', IMG)

# --- 2) va slide 8 ---
F='bao_cao/SLIDE_THUYET_TRINH_v6.pptx'
prs=Presentation(F); slide=prs.slides[7]

# thu font text 1.3 cho gon (de chua anh ben duoi)
for sh in slide.shapes:
    if sh.has_text_frame and sh.text_frame.text.startswith('1.3'):
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                cur = run.font.size
                run.font.size = Pt(12.5)
        # co lai chieu cao box
        sh.height = Inches(2.05); sh.top = Inches(1.98)

# chen anh phan bo o vung trong duoi-trai
pic = slide.shapes.add_picture(IMG, Inches(0.72), Inches(4.20), width=Inches(5.75))
# caption nho
cap = slide.shapes.add_textbox(Inches(0.72), Inches(6.02), Inches(5.75), Inches(0.28))
r = cap.text_frame.paragraphs[0].add_run()
r.text = 'Độ mòn dao tách rõ (hỏng ~177 vs lành ~123); tốc độ/mômen máy hỏng dồn về HAI ĐẦU.'
r.font.size = Pt(10); r.font.italic=True
from pptx.dml.color import RGBColor
r.font.color.rgb = RGBColor(0x33,0x33,0x33)

prs.save(F)
print('Slide 8 patched.')

# verify geometry
from pptx.util import Emu
p2=Presentation(F); s=p2.slides[7]
for sh in s.shapes:
    if sh.left is None: continue
    L,T,W,H=Emu(sh.left).inches,Emu(sh.top).inches,Emu(sh.width).inches,Emu(sh.height).inches
    tag='PIC' if sh.shape_type==13 else (repr(sh.text_frame.text[:32]) if sh.has_text_frame and sh.text_frame.text.strip() else sh.name)
    print(f'  L={L:.2f} T={T:.2f} R={L+W:.2f} B={T+H:.2f}  {tag}')
